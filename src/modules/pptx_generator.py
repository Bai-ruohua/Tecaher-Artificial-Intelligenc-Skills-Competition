# -*- coding: utf-8 -*-
"""
课件 PPTX 生成模块 v4.5
根据 AI 生成的课件大纲，创建 PowerPoint 幻灯片。
输出：data/uploads/<course_id>/<课程名>-课件.pptx
"""
import os
import re
from config import Config
from database import get_course_by_id, get_course_chapter
from deepseek_client import deepseek


def _course_name(course_id):
    try:
        c = get_course_by_id(course_id)
        return c["course_name"] if c else "本课程"
    except Exception:
        return "本课程"


def generate_courseware_outline(chapter_id, course_id, lesson_plan_text=""):
    """生成课件大纲（AI 调用），返回结构化文本"""
    ch = get_course_chapter(course_id, chapter_id)
    if not ch:
        return None

    cname = _course_name(course_id)
    sections = ch.get('key_points', ['主要内容'])
    prompt = f"""你是一位 PPT 课件设计专家。请为课程「{cname}」的章节「{ch['title']}」生成 PPT 逐页大纲。

教学目标：{'；'.join(ch.get('objectives', []))}
重点内容：{'；'.join(sections)}
{('教案内容参考：'+lesson_plan_text[:2000]) if lesson_plan_text else ''}

请生成 12-15 页的 PPT 大纲，按以下格式（JSON 数组）输出，不要有其他文字：

[
  {{"page":1,"title":"{cname} - {ch['title']}","subtitle":"授课教师：朱景峰","content":["宜宾工业职业技术学院 数字经济学院"]}},
  {{"page":2,"title":"学习目标","content":["目标1","目标2","目标3"]}},
  {{"page":3,"title":"知识点标题","content":["要点1","要点2","要点3"]}},
  ...
  {{"page":14,"title":"本节小结","content":["总结1","总结2","总结3"]}},
  {{"page":15,"title":"课后作业","content":["作业1","作业2"]}}
]

要求：
- 第1页为标题页
- 第2页为学习目标
- 第3-12页为知识点讲解（含代码示例或图示建议）
- 第13页为课堂练习
- 第14页为小结
- 第15页为课后作业
- 每页 content 2-5 条要点
- 知识点内容要紧扣章节重点，内容准确"""

    import json
    try:
        result = deepseek.chat("你是一个严谨的 JSON 数据输出助手。", prompt, max_tokens=4096)
        result = result.strip()
        if result.startswith("```"):
            result = result.split("\n", 1)[1] if "\n" in result else result[3:]
            if result.endswith("```"):
                result = result[:-3]
            result = result.strip()
        slides = json.loads(result)
        if isinstance(slides, list) and len(slides) > 0:
            return slides
    except Exception as e:
        print(f"  [pptx_gen] AI 生成失败: {e}")

    # 降级方案
    return _fallback_slides(cname, ch)


def _fallback_slides(cname, ch):
    """AI 失败时的降级幻灯片结构"""
    sections = ch.get('key_points', ['主要内容'])
    slides = [
        {"page": 1, "title": f"{cname} - {ch['title']}", "subtitle": "授课教师：朱景峰", "content": ["宜宾工业职业技术学院 数字经济学院"]},
        {"page": 2, "title": "学习目标", "content": ch.get('objectives', ['掌握核心概念'])},
    ]
    for i, section in enumerate(sections[:8]):
        slides.append({"page": 3+i, "title": section, "content": [f"{section}要点1", f"{section}要点2", f"{section}要点3"]})
    slides.append({"page": len(slides)+1, "title": "课堂练习", "content": ["完成课堂编程练习", "小组讨论与分享"]})
    slides.append({"page": len(slides)+1, "title": "本节小结", "content": ["梳理本章知识体系", "重点回顾", "常见错误提示"]})
    slides.append({"page": len(slides)+1, "title": "课后作业", "content": ["完成课后习题", "预习下一章节"]})
    return slides


def generate_pptx(chapter_id, course_id, lesson_plan_text=""):
    """
    生成课件 PPTX 文件。
    返回 {success, file_path, download_url, slide_count}
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    cname = _course_name(course_id)
    ch = get_course_chapter(course_id, chapter_id)
    if not ch:
        return {"success": False, "error": "章节不存在"}

    # 生成大纲
    slides_data = generate_courseware_outline(chapter_id, course_id, lesson_plan_text)
    if not slides_data:
        return {"success": False, "error": "大纲生成失败"}

    # 创建 PPT
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 主题色
    PRIMARY = RGBColor(0x1E, 0x40, 0xAF)      # 深蓝
    SECONDARY = RGBColor(0x39, 0x48, 0x9B)     # 中蓝
    ACCENT = RGBColor(0x3B, 0x82, 0xF6)        # 亮蓝
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x1F, 0x29, 0x37)
    GRAY = RGBColor(0x6B, 0x72, 0x80)

    for idx, slide_info in enumerate(slides_data):
        page = slide_info.get("page", idx+1)
        title = slide_info.get("title", "")
        subtitle = slide_info.get("subtitle", "")
        content = slide_info.get("content", [])

        if page == 1:
            # 标题页
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
            # 背景色块
            bg = slide.shapes.add_shape(
                1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = PRIMARY

            # 标题
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(40)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            if subtitle:
                txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
                tf2 = txBox2.text_frame
                p2 = tf2.paragraphs[0]
                p2.text = subtitle
                p2.font.size = Pt(24)
                p2.font.color.rgb = RGBColor(0xBB, 0xC7, 0xE6)
                p2.alignment = PP_ALIGN.CENTER

            # 底部学校信息
            txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.8))
            tf3 = txBox3.text_frame
            for ci, text in enumerate(content):
                if ci == 0:
                    p3 = tf3.paragraphs[0]
                else:
                    p3 = tf3.add_paragraph()
                p3.text = text
                p3.font.size = Pt(16)
                p3.font.color.rgb = RGBColor(0xBB, 0xC7, 0xE6)
                p3.alignment = PP_ALIGN.CENTER

        else:
            # 内容页
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 顶部色条
            bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
            bar.fill.solid()
            bar.fill.fore_color.rgb = ACCENT

            # 标题区域
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = PRIMARY

            # 内容区域
            txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for ci, item in enumerate(content):
                if ci == 0:
                    p2 = tf2.paragraphs[0]
                else:
                    p2 = tf2.add_paragraph()
                # 区分标题和子项
                if item.startswith("代码") or item.startswith("示例"):
                    p2.text = f"  {item}"
                    p2.font.size = Pt(16)
                    p2.font.color.rgb = ACCENT
                    p2.font.italic = True
                else:
                    p2.text = f"• {item}"
                    p2.font.size = Pt(18)
                    p2.font.color.rgb = DARK
                p2.space_after = Pt(8)

            # 页码
            txBox3 = slide.shapes.add_textbox(Inches(11.5), Inches(7), Inches(1.5), Inches(0.4))
            tf3 = txBox3.text_frame
            p3 = tf3.paragraphs[0]
            p3.text = str(page)
            p3.font.size = Pt(12)
            p3.font.color.rgb = GRAY
            p3.alignment = PP_ALIGN.RIGHT

    # 保存
    safe_name = re.sub(r'[\\/:*?"<>|]', '_', cname)
    out_dir = os.path.join(Config.UPLOAD_DIR, str(course_id))
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"{safe_name}-课件.pptx")

    if os.path.exists(out_path):
        base = f"{safe_name}-课件"
        n = 1
        while os.path.exists(os.path.join(out_dir, f"{base}_{n}.pptx")):
            n += 1
        out_path = os.path.join(out_dir, f"{base}_{n}.pptx")

    prs.save(out_path)

    download_url = f"/api/course/{course_id}/download/{os.path.basename(out_path)}"
    return {
        "success": True,
        "file_path": out_path,
        "download_url": download_url,
        "slide_count": len(slides_data)
    }
